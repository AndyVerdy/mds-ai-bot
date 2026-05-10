"""
Document ingestion pipeline.
Supports: .txt, .md, .vtt, .srt, .pdf, .pptx, plus WhatsApp digests from Airtable.
Chunks documents, embeds them, and stores in ChromaDB.

Improvements:
- Extracts speaker name from filename (e.g., "1. Josh Hadley_otter_ai.txt" → "Josh Hadley")
- Prepends speaker/source context header to each chunk
- Increased chunk size for better conversational context
- WhatsApp digests pulled from Airtable Summaries.raw_log (full conversation text)
"""

import json
import os
import re
import sys
from pathlib import Path

import fitz  # pymupdf
import requests
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma
from langchain_core.documents import Document
from rich.console import Console
from rich.progress import Progress

import config

# Airtable for WhatsApp digests
AIRTABLE_BASE_ID = "appT9TVZWhv7io4CN"
AIRTABLE_DIGESTS_TABLE = "Summaries"

console = Console()

# File metadata cache (loaded from data/metadata.json)
_file_metadata: dict = {}

# Video links cache (loaded from data/video_links.json)
_video_links: dict = {}


def load_file_metadata(data_dir: str = None) -> dict:
    """Load file metadata (dates, events, topics) from metadata.json."""
    global _file_metadata
    search_dirs = [data_dir] if data_dir else [str(config.DATA_DIR), "data/", "data"]
    for d in search_dirs:
        meta_path = Path(d) / "metadata.json"
        if meta_path.exists():
            try:
                with open(meta_path, "r", encoding="utf-8") as f:
                    raw = json.load(f)
                _file_metadata = raw
                return _file_metadata
            except Exception:
                pass
    return {}


def load_video_links(data_dir: str = None) -> dict:
    """Load video links mapping from video_links.json."""
    global _video_links
    search_dirs = [data_dir] if data_dir else [str(config.DATA_DIR), "data/", "data"]
    for d in search_dirs:
        vl_path = Path(d) / "video_links.json"
        if vl_path.exists():
            try:
                with open(vl_path, "r", encoding="utf-8") as f:
                    _video_links = json.load(f)
                console.print(f"[blue]Loaded {len(_video_links)} video links[/blue]")
                return _video_links
            except Exception:
                pass
    return {}


def get_video_url(file_path: str) -> str:
    """Get video URL for a transcript file, if available."""
    filename = Path(file_path).name
    entry = _video_links.get(filename)
    if entry:
        return entry.get("video_url", "")
    return ""


def get_file_metadata(file_path: str) -> dict:
    """Get date/event/topic metadata for a specific file."""
    filename = Path(file_path).name
    files = _file_metadata.get("files", {})
    defaults = _file_metadata.get("default", {})

    if filename in files:
        meta = {**defaults, **files[filename]}
    else:
        meta = dict(defaults)
    return meta


def extract_speaker_name(file_path: str) -> str:
    """Extract speaker name from filename.

    Handles patterns like:
        "1. Josh Hadley_otter_ai.txt" → "Josh Hadley"
        "11. Ozlem Sengul_otter_ai.txt" → "Ozlem Sengul"
        "Alar Huul_otter_ai.txt" → "Alar Huul"
        "Matt Dreier_otter_ai.txt" → "Matt Dreier"
    """
    name = Path(file_path).stem  # remove extension

    # Remove _otter_ai suffix
    name = re.sub(r"_otter_ai$", "", name, flags=re.IGNORECASE)

    # Remove leading number + dot prefix (e.g., "1. ", "11. ")
    name = re.sub(r"^\d+\.\s*", "", name)

    return name.strip()


def make_context_header(metadata: dict) -> str:
    """Build a context header string to prepend to chunk content."""
    parts = []

    speaker = metadata.get("speaker")
    if speaker:
        parts.append(f"Speaker: {speaker}")

    event = metadata.get("event")
    if event:
        parts.append(f"Event: {event}")

    date = metadata.get("date")
    if date:
        parts.append(f"Date: {date}")

    topic = metadata.get("topic")
    if topic:
        parts.append(f"Topic: {topic}")

    source = metadata.get("source", "")
    if source:
        parts.append(f"Source: {Path(source).name}")

    doc_type = metadata.get("type", "")
    if doc_type:
        parts.append(f"Type: {doc_type}")

    if metadata.get("timestamp_start"):
        parts.append(f"Timestamp: {metadata['timestamp_start']}")
    if metadata.get("page"):
        parts.append(f"Page: {metadata['page']}")
    if metadata.get("slide"):
        parts.append(f"Slide: {metadata['slide']}")

    if parts:
        return "[" + " | ".join(parts) + "]\n"
    return ""


def parse_vtt_srt(text: str, file_path: str) -> list[Document]:
    """Parse VTT/SRT subtitle files into timestamped chunks."""
    lines = text.strip().split("\n")
    segments = []
    current_time = ""
    current_text = []
    file_meta = get_file_metadata(file_path)

    timestamp_pattern = re.compile(
        r"(\d{1,2}:)?\d{2}:\d{2}[.,]\d{3}\s*-->\s*(\d{1,2}:)?\d{2}:\d{2}[.,]\d{3}"
    )

    for line in lines:
        line = line.strip()
        if timestamp_pattern.search(line):
            if current_text:
                segments.append({
                    "timestamp": current_time,
                    "text": " ".join(current_text),
                })
            current_time = line.split("-->")[0].strip()
            current_text = []
        elif line and not line.isdigit() and not line.startswith("WEBVTT"):
            clean = re.sub(r"<[^>]+>", "", line)
            if clean.strip():
                current_text.append(clean.strip())

    if current_text:
        segments.append({
            "timestamp": current_time,
            "text": " ".join(current_text),
        })

    speaker = extract_speaker_name(file_path)

    # Group segments into larger chunks
    docs = []
    chunk_text = []
    chunk_start = segments[0]["timestamp"] if segments else ""
    char_count = 0

    for seg in segments:
        chunk_text.append(seg["text"])
        char_count += len(seg["text"])
        if char_count >= config.CHUNK_SIZE:
            meta = {
                "source": file_path,
                "type": "transcript",
                "timestamp_start": chunk_start,
                "speaker": speaker,
                **{k: v for k, v in file_meta.items() if k != "_comment"},
            }
            content = make_context_header(meta) + " ".join(chunk_text)
            docs.append(Document(page_content=content, metadata=meta))
            chunk_text = []
            chunk_start = seg["timestamp"]
            char_count = 0

    if chunk_text:
        meta = {
            "source": file_path,
            "type": "transcript",
            "timestamp_start": chunk_start,
            "speaker": speaker,
            **{k: v for k, v in file_meta.items() if k != "_comment"},
        }
        content = make_context_header(meta) + " ".join(chunk_text)
        docs.append(Document(page_content=content, metadata=meta))

    return docs


def parse_otter_transcript(text: str, file_path: str) -> list[Document]:
    """Parse Otter.ai transcript .txt files with speaker/timestamp awareness.

    Format: "Unknown Speaker  0:05\ntext...\n\nUnknown Speaker  1:23\ntext..."
    """
    speaker_from_filename = extract_speaker_name(file_path)
    file_meta = get_file_metadata(file_path)

    # Split into speaker segments
    # Pattern: "Speaker Name  H:MM:SS" or "Unknown Speaker  M:SS"
    segment_pattern = re.compile(
        r"^(.+?)\s{2,}(\d{1,2}:\d{2}(?::\d{2})?)\s*$", re.MULTILINE
    )

    segments = []
    matches = list(segment_pattern.finditer(text))

    for i, match in enumerate(matches):
        seg_speaker = match.group(1).strip()
        timestamp = match.group(2).strip()
        start = match.end()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        content = text[start:end].strip()

        if content:
            segments.append({
                "speaker": seg_speaker,
                "timestamp": timestamp,
                "text": content,
            })

    if not segments:
        # Fallback: treat as plain text if no speaker patterns found
        meta = {
            "source": file_path,
            "type": "transcript",
            "speaker": speaker_from_filename,
            **{k: v for k, v in file_meta.items() if k != "_comment"},
        }
        return [Document(page_content=text, metadata=meta)]

    # Group segments into chunks of ~CHUNK_SIZE characters
    docs = []
    chunk_parts = []
    chunk_start_ts = segments[0]["timestamp"] if segments else ""
    char_count = 0

    for seg in segments:
        chunk_parts.append(seg["text"])
        char_count += len(seg["text"])

        if char_count >= config.CHUNK_SIZE:
            meta = {
                "source": file_path,
                "type": "transcript",
                "timestamp_start": chunk_start_ts,
                "speaker": speaker_from_filename,
                **{k: v for k, v in file_meta.items() if k != "_comment"},
            }
            content = make_context_header(meta) + "\n".join(chunk_parts)
            docs.append(Document(page_content=content, metadata=meta))
            chunk_parts = []
            chunk_start_ts = seg["timestamp"]
            char_count = 0

    if chunk_parts:
        meta = {
            "source": file_path,
            "type": "transcript",
            "timestamp_start": chunk_start_ts,
            "speaker": speaker_from_filename,
            **{k: v for k, v in file_meta.items() if k != "_comment"},
        }
        content = make_context_header(meta) + "\n".join(chunk_parts)
        docs.append(Document(page_content=content, metadata=meta))

    return docs


def parse_pdf(file_path: str) -> list[Document]:
    """Extract text from PDF, one chunk per page (then split further)."""
    doc = fitz.open(file_path)
    speaker = extract_speaker_name(file_path)
    file_meta = get_file_metadata(file_path)
    documents = []
    for page_num in range(len(doc)):
        page = doc[page_num]
        text = page.get_text().strip()
        if text:
            meta = {
                "source": file_path,
                "type": "pdf",
                "page": page_num + 1,
                "speaker": speaker,
                **{k: v for k, v in file_meta.items() if k != "_comment"},
            }
            content = make_context_header(meta) + text
            documents.append(Document(page_content=content, metadata=meta))
    doc.close()
    return documents


def parse_pptx(file_path: str) -> list[Document]:
    """Extract text from PPTX, one document per slide."""
    prs = Presentation(file_path)
    speaker = extract_speaker_name(file_path)
    file_meta = get_file_metadata(file_path)
    documents = []
    for slide_num, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            meta = {
                "source": file_path,
                "type": "presentation",
                "slide": slide_num,
                "speaker": speaker,
                **{k: v for k, v in file_meta.items() if k != "_comment"},
            }
            content = make_context_header(meta) + "\n".join(texts)
            documents.append(Document(page_content=content, metadata=meta))
    return documents


def parse_text(file_path: str) -> list[Document]:
    """Parse plain text or markdown files."""
    text = Path(file_path).read_text(encoding="utf-8")
    speaker = extract_speaker_name(file_path)
    file_meta = get_file_metadata(file_path)
    meta = {
        "source": file_path,
        "type": "text",
        "speaker": speaker,
        **{k: v for k, v in file_meta.items() if k != "_comment"},
    }
    content = make_context_header(meta) + text
    return [Document(page_content=content, metadata=meta)]


def is_otter_transcript(file_path: str) -> bool:
    """Detect if a .txt file is an Otter.ai transcript."""
    name = Path(file_path).name.lower()
    if "_otter_ai" in name:
        return True
    # Check first few lines for Otter format
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            head = f.read(500)
        return bool(re.search(r"Unknown Speaker\s{2,}\d+:\d{2}", head))
    except Exception:
        return False


def load_document(file_path: str) -> list[Document]:
    """Load and parse a document based on its extension."""
    path = Path(file_path)
    ext = path.suffix.lower()

    if ext in (".vtt", ".srt"):
        text = path.read_text(encoding="utf-8")
        return parse_vtt_srt(text, str(path))
    elif ext == ".pdf":
        return parse_pdf(str(path))
    elif ext == ".pptx":
        return parse_pptx(str(path))
    elif ext == ".txt":
        if is_otter_transcript(str(path)):
            text = path.read_text(encoding="utf-8")
            return parse_otter_transcript(text, str(path))
        return parse_text(str(path))
    elif ext == ".md":
        return parse_text(str(path))
    else:
        console.print(f"[yellow]Skipping unsupported file: {path.name}[/yellow]")
        return []


def chunk_documents(documents: list[Document]) -> list[Document]:
    """Split documents into smaller chunks for embedding."""
    # Transcripts are already chunked by timestamp/speaker
    already_chunked = [d for d in documents if d.metadata.get("type") == "transcript"]
    needs_chunking = [d for d in documents if d.metadata.get("type") != "transcript"]

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=config.CHUNK_SIZE,
        chunk_overlap=config.CHUNK_OVERLAP,
        separators=["\n\n", "\n", ". ", " ", ""],
    )

    chunked = splitter.split_documents(needs_chunking)
    return already_chunked + chunked


def ingest_files(file_paths: list[str]) -> int:
    """Ingest a list of files into the vector store. Returns number of chunks indexed."""
    all_docs = []
    with Progress() as progress:
        task = progress.add_task("Loading documents...", total=len(file_paths))
        for fp in file_paths:
            docs = load_document(fp)
            all_docs.extend(docs)
            progress.update(task, advance=1)

    if not all_docs:
        console.print("[yellow]No documents to ingest.[/yellow]")
        return 0

    console.print(f"Loaded {len(all_docs)} raw segments from {len(file_paths)} files")

    # Enrich documents with video URLs
    video_count = 0
    for doc in all_docs:
        source = doc.metadata.get("source", "")
        video_url = get_video_url(source)
        if video_url:
            doc.metadata["video_url"] = video_url
            video_count += 1
    if video_count:
        console.print(f"[blue]Attached video URLs to {video_count} segments[/blue]")

    # Chunk
    chunks = chunk_documents(all_docs)
    console.print(f"Split into {len(chunks)} chunks")

    # Embed and store using ChromaDB's built-in embeddings (free, local)
    console.print("Embedding and indexing (local model)...")
    vectorstore = Chroma(
        collection_name=config.COLLECTION_NAME,
        persist_directory=str(config.VECTORSTORE_DIR),
    )

    # Add in batches
    batch_size = 100
    with Progress() as progress:
        task = progress.add_task("Indexing chunks...", total=len(chunks))
        for i in range(0, len(chunks), batch_size):
            batch = chunks[i:i + batch_size]
            vectorstore.add_documents(batch)
            progress.update(task, advance=len(batch))

    console.print(f"[green]Indexed {len(chunks)} chunks into vector store.[/green]")
    return len(chunks)


def fetch_members_phone_to_name() -> dict[str, str]:
    """Build a {phone: full_name} map from the Airtable Members table.

    Used by the WhatsApp ingestion to enrich `@PushName` mentions with the
    member's full name (push names are often first-name-only or inconsistent
    across chats — the Members table has the canonical full name).

    Phones are stored without a leading `+` in both Whapi messages and the
    Members table, so a direct string match works.

    Returns an empty dict if AIRTABLE_PAT is missing or the call fails — the
    caller is expected to skip enrichment and proceed.
    """
    pat = os.getenv("AIRTABLE_PAT")
    if not pat:
        return {}
    url = f"https://api.airtable.com/v0/{AIRTABLE_BASE_ID}/Members"
    headers = {"Authorization": f"Bearer {pat}"}
    out: dict[str, str] = {}
    offset: str | None = None
    while True:
        params = {
            "pageSize": 100,
            "fields[]": ["phone", "name"],
        }
        if offset:
            params["offset"] = offset
        try:
            resp = requests.get(url, headers=headers, params=params, timeout=30)
            resp.raise_for_status()
        except Exception as e:
            console.print(
                f"[yellow]Members fetch failed (continuing without enrichment): {e}[/yellow]"
            )
            return out
        data = resp.json()
        for r in data.get("records", []):
            f = r.get("fields", {})
            phone = str(f.get("phone") or "").strip()
            name = (f.get("name") or "").strip()
            if phone and name:
                out[phone] = name
        offset = data.get("offset")
        if not offset:
            break
    console.print(f"[blue]Loaded {len(out)} Member phone -> name mappings[/blue]")
    return out


def build_pushname_to_fullname(
    source_messages_json: str,
    phone_to_name: dict[str, str],
) -> dict[str, str]:
    """For one digest, build {push_name: full_name} by matching each message's
    sender phone against the Members map. Returns empty when no matches."""
    if not source_messages_json or not phone_to_name:
        return {}
    try:
        msgs = json.loads(source_messages_json)
    except Exception:
        return {}
    out: dict[str, str] = {}
    for m in msgs:
        push = (m.get("from_name") or "").strip()
        phone = str(m.get("from") or "").strip()
        if not push or not phone or push in out:
            continue
        full = phone_to_name.get(phone)
        if full and full != push:
            out[push] = full
    return out


def enrich_raw_log(raw_log: str, pushname_to_fullname: dict[str, str]) -> str:
    """Substitute `@PushName:` with `@FullName (PushName):` in the raw log.

    Keeps the push name in parens so the embedded text still matches the
    original Whapi handle — useful when Members has stale full names but the
    push name is what other docs (and the Claude attribution prompt) use.
    """
    if not pushname_to_fullname:
        return raw_log
    out = raw_log
    # Longest first to avoid prefix collisions (e.g. "Aaron" vs "Aaron K").
    for push in sorted(pushname_to_fullname.keys(), key=len, reverse=True):
        full = pushname_to_fullname[push]
        out = out.replace(f"@{push}:", f"@{full} ({push}):")
    return out


def fetch_whatsapp_digests(min_msg_count: int = 1) -> list[dict]:
    """Fetch all daily WhatsApp digests from Airtable Summaries that have a
    non-empty raw_log and at least `min_msg_count` messages.

    Returns a list of digest dicts (Airtable fields with `id` added).
    Requires AIRTABLE_PAT in the environment.
    """
    pat = os.getenv("AIRTABLE_PAT")
    if not pat:
        console.print("[yellow]AIRTABLE_PAT not set — skipping WhatsApp ingestion[/yellow]")
        return []

    url = f"https://api.airtable.com/v0/{AIRTABLE_BASE_ID}/{AIRTABLE_DIGESTS_TABLE}"
    headers = {"Authorization": f"Bearer {pat}"}
    formula = f"AND({{period_type}}='daily',{{msg_count}}>={min_msg_count},{{raw_log}}!='')"

    digests: list[dict] = []
    offset: str | None = None
    page_count = 0
    while True:
        params = {
            "pageSize": 100,
            "sort[0][field]": "date",
            "sort[0][direction]": "desc",
            "filterByFormula": formula,
        }
        if offset:
            params["offset"] = offset
        resp = requests.get(url, headers=headers, params=params, timeout=30)
        resp.raise_for_status()
        data = resp.json()
        for r in data.get("records", []):
            f = r.get("fields", {})
            digests.append({
                "id": r["id"],
                "date": f.get("date"),
                "chat_id": f.get("chat_id"),
                "chat_name": f.get("chat_name"),
                "period_type": f.get("period_type"),
                "msg_count": f.get("msg_count", 0) or 0,
                "participant_count": f.get("participant_count", 0) or 0,
                "raw_log": f.get("raw_log", "") or "",
                "topics": f.get("topics", "") or "",
                "tl_dr": f.get("tl_dr", "") or "",
                # Pulled so make_whatsapp_documents can join sender phones to
                # the Members table for full-name enrichment (#12).
                "source_messages_json": f.get("source_messages_json", "") or "",
            })
        offset = data.get("offset")
        page_count += 1
        if not offset:
            break
    console.print(f"[blue]Fetched {len(digests)} WhatsApp digests across {page_count} pages[/blue]")
    return digests


def make_whatsapp_documents(digests: list[dict]) -> list[Document]:
    """Convert WhatsApp digests into LangChain Documents, chunked by character size."""
    if not digests:
        return []

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=config.CHUNK_SIZE,
        chunk_overlap=config.CHUNK_OVERLAP,
        # Prefer to break between messages (each starts with `[YYYY-MM-DD HH:MM UTC]`).
        separators=["\n[", "\n\n", "\n", ". ", " ", ""],
    )

    # Build the Members phone -> full-name map ONCE for the whole batch (#12).
    # Then enrich each digest's raw_log so @PushName becomes @FullName (PushName).
    phone_to_name = fetch_members_phone_to_name()

    enriched_count = 0
    docs: list[Document] = []
    for d in digests:
        raw = d.get("raw_log") or ""
        if not raw.strip():
            continue

        # Enrich raw_log with full names from Members (#12).
        pushname_map = build_pushname_to_fullname(
            d.get("source_messages_json", ""), phone_to_name
        )
        if pushname_map:
            raw = enrich_raw_log(raw, pushname_map)
            enriched_count += len(pushname_map)

        chat_name = d.get("chat_name") or "Unknown group"
        date = d.get("date") or ""
        period = d.get("period_type") or "daily"
        chat_id = d.get("chat_id") or ""
        tl_dr = (d.get("tl_dr") or "").strip()
        topics = (d.get("topics") or "").strip()

        # Embedding-friendly header: includes tl_dr + topics so semantic search
        # has more signal even before the raw chat log starts.
        header_lines = [
            f"[Source: WhatsApp conversation | Group: {chat_name} | "
            f"Date: {date} | Period: {period}]"
        ]
        if tl_dr:
            header_lines.append(f"Summary: {tl_dr}")
        if topics:
            header_lines.append(f"Topics: {topics}")
        header = "\n".join(header_lines)
        # Single document per digest, will be split below.
        meta = {
            "source": f"airtable://Summaries/{d['id']}",
            "source_id": d["id"],
            "type": "whatsapp",
            "chat_name": chat_name,
            "chat_id": chat_id,
            "date": date,
            "period_type": period,
            "msg_count": int(d.get("msg_count") or 0),
            "participant_count": int(d.get("participant_count") or 0),
            # Leave speaker empty — WhatsApp digests are multi-participant.
            "speaker": "",
        }
        full_doc = Document(page_content=f"{header}\n{raw}", metadata=meta)
        # Split if needed; each chunk gets its own header so retrieval still has context.
        if len(full_doc.page_content) <= config.CHUNK_SIZE:
            docs.append(full_doc)
        else:
            for chunk in splitter.split_documents([full_doc]):
                # Re-prepend the header to every chunk (split may have removed it).
                if not chunk.page_content.lstrip().startswith("[Source: WhatsApp"):
                    chunk = Document(
                        page_content=f"{header}\n{chunk.page_content}",
                        metadata=chunk.metadata,
                    )
                docs.append(chunk)

    console.print(
        f"[blue]Built {len(docs)} WhatsApp document chunks from {len(digests)} digests "
        f"(enriched {enriched_count} push-name -> full-name mappings)[/blue]"
    )
    return docs


def ingest_whatsapp(min_msg_count: int = 1) -> int:
    """Pull all qualifying WhatsApp digests from Airtable, chunk, embed, store.
    Returns the number of chunks added."""
    digests = fetch_whatsapp_digests(min_msg_count=min_msg_count)
    if not digests:
        return 0
    chunks = make_whatsapp_documents(digests)
    if not chunks:
        return 0

    console.print("Embedding and indexing WhatsApp chunks (local model)...")
    vectorstore = Chroma(
        collection_name=config.COLLECTION_NAME,
        persist_directory=str(config.VECTORSTORE_DIR),
    )

    batch_size = 100
    with Progress() as progress:
        task = progress.add_task("Indexing WA chunks...", total=len(chunks))
        for i in range(0, len(chunks), batch_size):
            batch = chunks[i:i + batch_size]
            vectorstore.add_documents(batch)
            progress.update(task, advance=len(batch))

    console.print(f"[green]Indexed {len(chunks)} WhatsApp chunks into vector store.[/green]")
    return len(chunks)


def make_video_documents(video: dict, segments: list[dict]) -> list[Document]:
    """Group transcript segments into chunks bounded by chapter and ~1800 chars.

    Walks segments in start_ms order. Starts a new chunk when:
      - the chapter changes, OR
      - the current chunk's body is approaching the chunk-size cap.

    Each chunk gets a context header (`[Video: … · Chapter: … · Speaker: …]`)
    prepended to the body so the LLM sees the source structure even when only
    one chunk is retrieved.

    Metadata on each Document carries the deep-link primitives (`video_id`,
    `start_ms`, `chapter_title`) so search results can render with a "jump to
    this chapter" affordance.
    """
    if not segments:
        return []

    docs: list[Document] = []
    current_chunk: list[dict] = []
    current_chapter: str = (segments[0].get("chapter_title") or "")
    current_chars: int = 0
    chunk_char_cap = 1800

    for seg in segments:
        seg_chapter = seg.get("chapter_title") or ""
        seg_text = seg.get("text") or ""

        if current_chunk and (
            seg_chapter != current_chapter
            or current_chars + len(seg_text) > chunk_char_cap
        ):
            doc = _build_video_doc(video, current_chunk, current_chapter)
            if doc is not None:
                docs.append(doc)
            current_chunk = []
            current_chars = 0
            current_chapter = seg_chapter

        current_chunk.append(seg)
        current_chars += len(seg_text)

    if current_chunk:
        doc = _build_video_doc(video, current_chunk, current_chapter)
        if doc is not None:
            docs.append(doc)

    return docs


def _build_video_doc(video: dict, segs: list[dict], chapter: str):
    """Build a single ChromaDB Document for a chunk of consecutive segments."""
    if not segs:
        return None

    first = segs[0]
    last = segs[-1]
    speaker_label = first.get("speaker_label") or ""
    start_ms = int(first.get("start_ms") or 0)
    end_ms = int(last.get("end_ms") or 0)

    body = " ".join((s.get("text") or "").strip() for s in segs if s.get("text"))
    if not body.strip():
        return None

    header_parts = [f"Video: {video.get('title') or 'Untitled'}"]
    if chapter:
        header_parts.append(f"Chapter: {chapter}")
    if speaker_label:
        header_parts.append(f"Speaker: {speaker_label}")
    recorded = video.get("recorded_at") or video.get("uploaded_at") or ""
    if recorded:
        header_parts.append(f"Recorded: {recorded[:10]}")
    header = "[" + " · ".join(header_parts) + "]"

    page_content = f"{header}\n{body}"

    metadata = {
        "type": "video",
        "source": f"video://{video['id']}#t={start_ms // 1000}",
        "source_id": video["id"],
        "video_id": video["id"],
        "video_title": video.get("title") or "",
        "thumbnail_url": video.get("thumbnail_url") or "",
        "chapter_title": chapter or "",
        "speaker_label": speaker_label,
        "start_ms": start_ms,
        "end_ms": end_ms,
        "duration_sec": int(video.get("duration_sec") or 0),
        "date": (recorded or "")[:10],
        # Backwards-compat: existing source enrichment looks at `speaker`
        # for display when no specific source_type branch matched. Setting
        # it to the video title gives any naive consumer something to render.
        "speaker": video.get("title") or "",
    }

    return Document(page_content=page_content, metadata=metadata)


def ingest_videos(force: bool = False) -> int:
    """Pull videos with ready transcripts from Postgres, chunk transcript
    segments, embed, and store in ChromaDB. Returns count of chunks added.

    Reads from:
      - `videos` (mux_status='ready' AND transcription_status='ready' AND deleted_at IS NULL)
      - `transcript_segments` (joined per-video, ordered by start_ms)

    Mirrors `ingest_whatsapp` — same vector store, same embedding model. Search
    via `/api/ask` automatically ranks video chunks alongside WA + legacy.
    """
    from videos import _supabase_get

    videos = _supabase_get(
        "videos",
        {
            "select": "id,title,thumbnail_url,duration_sec,recorded_at,uploaded_at",
            "mux_status": "eq.ready",
            "transcription_status": "eq.ready",
            "deleted_at": "is.null",
            "order": "uploaded_at.desc",
        },
    )

    if not videos:
        console.print("[yellow]No ready videos found.[/yellow]")
        return 0

    console.print(f"Found {len(videos)} ready videos")

    all_chunks: list[Document] = []
    for v in videos:
        segs = _supabase_get(
            "transcript_segments",
            {
                "select": "text,start_ms,end_ms,speaker_label,chapter_title",
                "video_id": f"eq.{v['id']}",
                "order": "start_ms.asc",
            },
        )
        if not segs:
            console.print(
                f"[yellow]Video {v['id'][:8]} ({v.get('title', '?')}): no segments[/yellow]"
            )
            continue
        chunks = make_video_documents(v, segs)
        title_short = (v.get("title") or "")[:40]
        console.print(
            f"Video {v['id'][:8]} ({title_short}): {len(segs)} segments → {len(chunks)} chunks"
        )
        all_chunks.extend(chunks)

    if not all_chunks:
        console.print("[yellow]No video chunks to ingest.[/yellow]")
        return 0

    console.print(f"Embedding and indexing {len(all_chunks)} video chunks (local model)...")
    vectorstore = Chroma(
        collection_name=config.COLLECTION_NAME,
        persist_directory=str(config.VECTORSTORE_DIR),
    )

    batch_size = 100
    with Progress() as progress:
        task = progress.add_task("Indexing video chunks...", total=len(all_chunks))
        for i in range(0, len(all_chunks), batch_size):
            batch = all_chunks[i:i + batch_size]
            vectorstore.add_documents(batch)
            progress.update(task, advance=len(batch))

    console.print(f"[green]Indexed {len(all_chunks)} video chunks into vector store.[/green]")
    return len(all_chunks)


def ingest_directory(directory: str) -> int:
    """Ingest all supported files from a directory."""
    # Load file metadata (dates, events) before ingesting
    load_file_metadata(directory)
    # Load video links (transcript → video URL mapping)
    load_video_links(str(Path(directory).parent))

    path = Path(directory)
    supported = {".txt", ".md", ".vtt", ".srt", ".pdf", ".pptx"}
    files = [str(f) for f in path.rglob("*") if f.suffix.lower() in supported]

    if not files:
        console.print(f"[yellow]No supported files found in {directory}[/yellow]")
        return 0

    console.print(f"Found {len(files)} files in {directory}")
    return ingest_files(files)


if __name__ == "__main__":
    if len(sys.argv) < 2:
        console.print("Usage:")
        console.print("  python ingest.py <file_or_directory> [file2] ...   # files")
        console.print("  python ingest.py --whatsapp                          # WhatsApp digests from Airtable")
        console.print("Supported formats: .txt, .md, .vtt, .srt, .pdf, .pptx")
        sys.exit(1)

    args = sys.argv[1:]
    total_chunks = 0

    if "--whatsapp" in args or "--wa" in args:
        total_chunks += ingest_whatsapp()
        args = [a for a in args if a not in ("--whatsapp", "--wa")]

    if "--videos" in args or "--video" in args:
        total_chunks += ingest_videos()
        args = [a for a in args if a not in ("--videos", "--video")]

    for p in args:
        path = Path(p)
        if path.is_dir():
            total_chunks += ingest_directory(str(path))
        elif path.is_file():
            total_chunks += ingest_files([str(path)])
        else:
            console.print(f"[red]Not found: {p}[/red]")

    console.print(f"\n[bold green]Done! Total chunks indexed: {total_chunks}[/bold green]")
